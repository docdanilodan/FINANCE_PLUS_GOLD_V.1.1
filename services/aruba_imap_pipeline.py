from __future__ import annotations

import hashlib
import imaplib
import io
import os
import re
import ssl
from datetime import date, datetime
from email import policy
from email.parser import BytesParser
from email.utils import parseaddr, parsedate_to_datetime

from docx import Document
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseUpload
from openpyxl import load_workbook
from pypdf import PdfReader
from pptx import Presentation

from document_ai import classify_text, suggested_name
from services.airtable_adapter import AirtableGold
from services.client_practice_matcher import match_client_practice
from services.gmail_drive_pipeline import _airtable_type, _archive_destination
from services.google_auth import load_credentials

DEFAULT_IMAP_HOST = "imaps.aruba.it"
DEFAULT_IMAP_PORT = 993


def _connect(account_email: str, password: str, host: str = DEFAULT_IMAP_HOST, port: int = DEFAULT_IMAP_PORT):
    if not account_email or not password:
        raise ValueError("Email e password Aruba sono obbligatorie.")
    context = ssl.create_default_context()
    client = imaplib.IMAP4_SSL(host=host, port=int(port), ssl_context=context)
    client.login(account_email, password)
    return client


def test_aruba_connection(
    account_email: str,
    password: str,
    host: str = DEFAULT_IMAP_HOST,
    port: int = DEFAULT_IMAP_PORT,
) -> dict:
    client = _connect(account_email, password, host, port)
    try:
        status, data = client.select("INBOX", readonly=True)
        if status != "OK":
            raise RuntimeError("Impossibile aprire INBOX.")
        count = int((data or [b"0"])[0] or 0)
        return {"ok": True, "account": account_email, "inbox_messages": count}
    finally:
        try:
            client.logout()
        except Exception:
            pass


def _decode_text(part) -> str:
    try:
        payload = part.get_payload(decode=True) or b""
        charset = part.get_content_charset() or "utf-8"
        return payload.decode(charset, errors="replace")
    except Exception:
        return ""


def _message_body(msg) -> str:
    plain: list[str] = []
    html: list[str] = []
    if msg.is_multipart():
        for part in msg.walk():
            if part.get_content_disposition() == "attachment":
                continue
            if part.get_content_type() == "text/plain":
                plain.append(_decode_text(part))
            elif part.get_content_type() == "text/html":
                html.append(_decode_text(part))
    elif msg.get_content_type() == "text/plain":
        plain.append(_decode_text(msg))
    elif msg.get_content_type() == "text/html":
        html.append(_decode_text(msg))

    if plain:
        return "\n".join(plain).strip()
    if html:
        text = re.sub(r"<script\b[^>]*>.*?</script>", " ", "\n".join(html), flags=re.I | re.S)
        text = re.sub(r"<style\b[^>]*>.*?</style>", " ", text, flags=re.I | re.S)
        text = re.sub(r"<[^>]+>", " ", text)
        return re.sub(r"\s+", " ", text).strip()
    return ""


def _pdf_text(raw: bytes) -> str:
    reader = PdfReader(io.BytesIO(raw))
    return "\n".join((page.extract_text() or "") for page in reader.pages)


def _docx_text(raw: bytes) -> str:
    doc = Document(io.BytesIO(raw))
    chunks = [paragraph.text for paragraph in doc.paragraphs if paragraph.text]
    for table in doc.tables:
        for row in table.rows:
            chunks.append(" | ".join(cell.text for cell in row.cells))
    return "\n".join(chunks)


def _xlsx_text(raw: bytes) -> str:
    workbook = load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
    chunks: list[str] = []
    cells = 0
    try:
        for sheet in workbook.worksheets:
            chunks.append(f"FOGLIO: {sheet.title}")
            for row in sheet.iter_rows(values_only=True):
                values = [str(value) for value in row if value not in (None, "")]
                if values:
                    chunks.append(" | ".join(values))
                cells += len(row)
                if cells >= 12000:
                    return "\n".join(chunks)
    finally:
        workbook.close()
    return "\n".join(chunks)


def _pptx_text(raw: bytes) -> str:
    presentation = Presentation(io.BytesIO(raw))
    chunks: list[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                chunks.append(shape.text)
    return "\n".join(chunks)


def _attachment_text(filename: str, mime_type: str, raw: bytes) -> str:
    ext = os.path.splitext(filename or "")[1].lower()
    try:
        if ext == ".pdf" or mime_type == "application/pdf":
            text = _pdf_text(raw)
        elif ext == ".docx":
            text = _docx_text(raw)
        elif ext in {".xlsx", ".xlsm", ".xltx"}:
            text = _xlsx_text(raw)
        elif ext == ".pptx":
            text = _pptx_text(raw)
        elif ext in {".txt", ".csv", ".md", ".xml", ".json", ".html", ".htm"} or mime_type.startswith("text/"):
            text = raw.decode("utf-8", errors="replace")
        else:
            text = ""
    except Exception:
        return ""
    return text[:160000]


def _imap_since(value: str | date | None) -> str:
    if value is None:
        d = date.today()
    elif isinstance(value, date):
        d = value
    else:
        d = date.fromisoformat(str(value))
    return d.strftime("%d-%b-%Y")


def _message_datetime(msg) -> str | None:
    raw = str(msg.get("Date") or "").strip()
    if not raw:
        return None
    try:
        return parsedate_to_datetime(raw).isoformat()
    except Exception:
        return None


def _append_source(existing: str, account_email: str) -> str:
    sources = [x.strip() for x in str(existing or "").replace(";", "\n").splitlines() if x.strip()]
    if account_email.casefold() not in {x.casefold() for x in sources}:
        sources.append(account_email)
    return "\n".join(sources)


def _find_source_message(airtable: AirtableGold, source_key: str) -> dict | None:
    return (
        airtable.find_one("email", "Message ID sorgente", source_key)
        or airtable.find_one("email", "Gmail Message ID", source_key)
    )


def _existing_source_keys(airtable: AirtableGold, account_email: str) -> set[str]:
    keys: set[str] = set()
    try:
        records = airtable.list_records("email", max_records=8000)
    except Exception:
        return keys
    target = account_email.casefold()
    for record in records:
        fields = record.get("fields", {})
        mailbox = str(fields.get("Casella origine") or fields.get("Casella sorgente") or "").casefold()
        if mailbox != target:
            continue
        value = str(fields.get("Message ID sorgente") or "").strip()
        if value:
            keys.add(value)
    return keys


def _register_duplicate_source(airtable: AirtableGold, sha: str, account_email: str) -> bool:
    existing = airtable.find_one("documenti", "SHA-256", sha)
    if not existing:
        return False
    fields = existing.get("fields", {})
    merged = _append_source(fields.get("Caselle origine", ""), account_email)
    updates = {}
    if merged != str(fields.get("Caselle origine", "") or ""):
        updates["Caselle origine"] = merged
    if not fields.get("Casella sorgente"):
        updates["Casella sorgente"] = account_email
    if updates:
        airtable.update_record("documenti", existing["id"], updates)
    return True


def _date_is_newer(candidate: str | None, current: str | None) -> bool:
    if not candidate:
        return False
    if not current:
        return True
    try:
        return date.fromisoformat(candidate[:10]) > date.fromisoformat(str(current)[:10])
    except ValueError:
        return True


def _update_client_document_state(
    airtable: AirtableGold,
    client_id: str,
    category: str,
    document_year: int | None,
    reference_date: str | None,
    filename: str,
) -> bool:
    try:
        client = airtable.get_record("clienti", client_id)
    except Exception:
        return False
    fields = client.get("fields", {})
    updates = {}

    if category == "Bilancio d'esercizio" and document_year:
        current_year = fields.get("Ultimo bilancio disponibile")
        try:
            current_year = int(current_year) if current_year not in (None, "") else None
        except (TypeError, ValueError):
            current_year = None
        if current_year is None or int(document_year) > current_year:
            updates["Ultimo bilancio disponibile"] = int(document_year)

    if category == "Centrale Rischi Banca d'Italia" and reference_date:
        if _date_is_newer(reference_date, fields.get("CR aggiornata al")):
            updates["CR aggiornata al"] = reference_date[:10]

    if category == "Visura Camerale" and reference_date:
        if _date_is_newer(reference_date, fields.get("Data estrazione visura")):
            updates["Data estrazione visura"] = reference_date[:10]
            updates["File sorgente visura"] = filename

    if not updates:
        return False
    airtable.update_record("clienti", client_id, updates)
    return True


def sync_aruba_attachments(
    account_email: str,
    password: str,
    since: str | date | None = None,
    drive_folder_id: str | None = None,
    max_messages: int = 100,
    host: str = DEFAULT_IMAP_HOST,
    port: int = DEFAULT_IMAP_PORT,
) -> dict:
    """Read one Aruba mailbox and apply the FinancePlus document workflow.

    The flow mirrors the Gmail archive: content-aware classification, client/practice
    matching, global SHA-256 deduplication, Drive filing and Airtable indexing.
    On a backlog, messages are processed oldest-first; already indexed IMAP UIDs are
    skipped so scheduled runs automatically advance until the historical archive is done.
    """
    drive = build("drive", "v3", credentials=load_credentials(), cache_discovery=False)
    airtable = AirtableGold()
    client = _connect(account_email, password, host, port)
    existing_source_keys = _existing_source_keys(airtable, account_email)
    result = {
        "account": account_email,
        "messages": 0,
        "attachments": 0,
        "uploaded": 0,
        "duplicates": 0,
        "duplicate_sources_updated": 0,
        "matched": 0,
        "archived_by_client": 0,
        "pending_review": 0,
        "client_updates": 0,
        "skipped_indexed_messages": 0,
        "errors": [],
    }

    try:
        status, _ = client.select("INBOX", readonly=True)
        if status != "OK":
            raise RuntimeError("Impossibile aprire INBOX.")

        status, data = client.uid("search", None, f'(SINCE "{_imap_since(since)}")')
        if status != "OK":
            raise RuntimeError("Ricerca IMAP non riuscita.")
        uids = [u for u in (data[0] or b"").split() if u]
        processed_new = 0

        for uid in uids:
            uid_text = uid.decode("ascii", errors="ignore")
            source_key = f"ARUBA:{account_email.casefold()}:UID:{uid_text}"
            if source_key in existing_source_keys:
                result["skipped_indexed_messages"] += 1
                continue
            if max_messages > 0 and processed_new >= int(max_messages):
                break

            try:
                status, payload = client.uid("fetch", uid, "(RFC822)")
                if status != "OK" or not payload or not payload[0]:
                    raise RuntimeError("Messaggio IMAP non leggibile.")
                raw_message = payload[0][1]
                msg = BytesParser(policy=policy.default).parsebytes(raw_message)

                message_id = str(msg.get("Message-ID") or "").strip()
                legacy_key = f"ARUBA:{account_email.casefold()}:{message_id or uid_text}"
                if legacy_key != source_key and _find_source_message(airtable, legacy_key):
                    result["duplicates"] += 1
                    existing_source_keys.add(source_key)
                    continue

                processed_new += 1
                result["messages"] += 1
                subject = str(msg.get("Subject") or "(senza oggetto)")
                sender_raw = str(msg.get("From") or "")
                sender_email = parseaddr(sender_raw)[1] or sender_raw
                body = _message_body(msg)
                context = f"{subject}\n{sender_raw}\n{body[:16000]}"
                email_match = match_client_practice(airtable, context)
                email_confident = bool(email_match.client_id and email_match.confidence >= 0.80)
                attachment_names: list[str] = []

                for part in msg.iter_attachments():
                    filename = str(part.get_filename() or "").strip()
                    if not filename:
                        continue
                    raw = part.get_payload(decode=True) or b""
                    if not raw:
                        continue
                    attachment_names.append(filename)
                    result["attachments"] += 1
                    sha = hashlib.sha256(raw).hexdigest()

                    if _register_duplicate_source(airtable, sha, account_email):
                        result["duplicates"] += 1
                        result["duplicate_sources_updated"] += 1
                        continue

                    extracted = _attachment_text(filename, part.get_content_type(), raw)
                    attachment_context = f"{filename}\n{context}\n{extracted}"
                    file_match = match_client_practice(airtable, attachment_context[:160000])
                    confident_client = bool(file_match.client_id and file_match.confidence >= 0.80)
                    if confident_client:
                        result["matched"] += 1

                    classification = classify_text(attachment_context[:160000])
                    ext = os.path.splitext(filename)[1] or ".bin"
                    if confident_client:
                        classification.company_name = file_match.client_name or ""
                    proposed = suggested_name(classification, extension=ext)

                    destination_id, archive_path = _archive_destination(
                        drive,
                        drive_folder_id,
                        file_match.client_name if confident_client else None,
                        classification.category,
                    )
                    media = MediaIoBaseUpload(
                        io.BytesIO(raw),
                        mimetype=part.get_content_type() or "application/octet-stream",
                        resumable=False,
                    )
                    saved = drive.files().create(
                        body={"name": proposed or filename, "parents": [destination_id]},
                        media_body=media,
                        fields="id,webViewLink,name,parents",
                    ).execute()
                    result["uploaded"] += 1

                    fields = {
                        "Documento": saved["name"],
                        "Tipo Documento": _airtable_type(classification.category),
                        "Nome Originale": filename,
                        "Nome IA Suggerito": saved["name"],
                        "Nome Definitivo": saved["name"],
                        "Origine": "Altro",
                        "Caselle origine": account_email,
                        "Casella sorgente": account_email,
                        "URL Drive": saved.get("webViewLink", ""),
                        "SHA-256": sha,
                        "Stato Verifica": "Da verificare",
                        "Percorso nel pacchetto": f"EMAIL_ARUBA/{account_email}/{archive_path}",
                    }
                    if classification.document_year:
                        fields["Esercizio"] = int(classification.document_year)
                    if classification.reference_date:
                        fields["Data Documento"] = classification.reference_date[:10]
                    if extracted:
                        fields["Sintesi IA"] = re.sub(r"\s+", " ", extracted)[:1800]
                    if confident_client:
                        fields["Cliente"] = file_match.client_name or ""
                        fields["Cliente collegato"] = [file_match.client_id]
                        result["archived_by_client"] += 1
                    else:
                        result["pending_review"] += 1
                    if file_match.practice_id and confident_client:
                        fields["Pratica ID"] = file_match.practice_code or ""
                        fields["Pratica collegata"] = [file_match.practice_id]
                    airtable.create_record("documenti", fields)

                    if confident_client and _update_client_document_state(
                        airtable,
                        file_match.client_id,
                        classification.category,
                        classification.document_year,
                        classification.reference_date,
                        saved["name"],
                    ):
                        result["client_updates"] += 1

                snippet = body.replace("\r", " ").replace("\n", " ")[:1500]
                email_fields = {
                    "Oggetto": subject,
                    "Mittente": sender_email,
                    "Casella origine": account_email,
                    "Casella sorgente": account_email,
                    "Message ID sorgente": source_key,
                    "Sintesi IA": snippet,
                    "Allegati": "\n".join(attachment_names),
                }
                msg_dt = _message_datetime(msg)
                if msg_dt:
                    email_fields["Data e ora"] = msg_dt
                if email_confident:
                    email_fields["Cliente"] = email_match.client_name or ""
                    email_fields["Cliente collegato"] = [email_match.client_id]
                if email_match.practice_id and email_confident:
                    email_fields["Pratica ID"] = email_match.practice_code or ""
                    email_fields["Pratica collegata"] = [email_match.practice_id]
                airtable.create_record("email", email_fields)
                existing_source_keys.add(source_key)

            except Exception as exc:
                result["errors"].append({"uid": uid_text, "error": str(exc)})
    finally:
        try:
            client.logout()
        except Exception:
            pass

    return result
