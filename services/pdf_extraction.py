from __future__ import annotations

import io
import os
from dataclasses import dataclass, field
from functools import lru_cache

from pypdf import PdfReader


@dataclass
class ExtractionResult:
    text: str = ""
    method: str = "none"
    cloud_used: bool = False
    warnings: list[str] = field(default_factory=list)


def _truthy(name: str, default: bool = False) -> bool:
    raw = os.getenv(name)
    if raw is None:
        return default
    return raw.strip().lower() in {"1", "true", "yes", "on"}


def _local_pdf_text(raw: bytes, max_chars: int = 250_000) -> str:
    reader = PdfReader(io.BytesIO(raw))
    chunks: list[str] = []
    total = 0
    for page in reader.pages:
        text = page.extract_text() or ""
        if not text:
            continue
        remaining = max_chars - total
        if remaining <= 0:
            break
        chunks.append(text[:remaining])
        total += min(len(text), remaining)
    return "\n\n".join(chunks).strip()


def _ocr_languages() -> str:
    return os.getenv("FINANCEPLUS_OCR_LANGUAGES", "ita+eng").strip() or "ita+eng"


@lru_cache(maxsize=1)
def _available_ocr_languages() -> set[str]:
    import pytesseract

    return set(pytesseract.get_languages(config=""))


def _ocr_image(image) -> str:
    import pytesseract

    requested = _ocr_languages().split("+")
    available = _available_ocr_languages()
    selected = [language for language in requested if language in available]
    if not selected and "eng" in available:
        selected = ["eng"]
    language_arg = "+".join(selected) if selected else None
    return pytesseract.image_to_string(image, lang=language_arg).strip()


def _local_image_ocr(raw: bytes) -> str:
    from PIL import Image, ImageOps

    try:
        from pillow_heif import register_heif_opener

        register_heif_opener()
    except ImportError:
        pass

    with Image.open(io.BytesIO(raw)) as image:
        image = ImageOps.exif_transpose(image).convert("RGB")
        return _ocr_image(image)[:250_000]


def _local_scanned_pdf_ocr(raw: bytes, max_chars: int = 250_000) -> str:
    import pypdfium2 as pdfium

    max_pages = max(1, min(int(os.getenv("FINANCEPLUS_OCR_MAX_PAGES", "30")), 100))
    document = pdfium.PdfDocument(io.BytesIO(raw))
    chunks: list[str] = []
    total = 0
    try:
        for index in range(min(len(document), max_pages)):
            page = document[index]
            try:
                image = page.render(scale=2).to_pil()
                text = _ocr_image(image)
            finally:
                page.close()
            if not text:
                continue
            remaining = max_chars - total
            if remaining <= 0:
                break
            chunks.append(text[:remaining])
            total += min(len(text), remaining)
    finally:
        document.close()
    return "\n\n".join(chunks).strip()


def _local_office_text(
    raw: bytes, lower_name: str, max_chars: int = 250_000
) -> tuple[str, str]:
    if lower_name.endswith(".docx"):
        from docx import Document

        doc = Document(io.BytesIO(raw))
        chunks = [paragraph.text for paragraph in doc.paragraphs if paragraph.text]
        for table in doc.tables:
            chunks.extend(
                " | ".join(cell.text for cell in row.cells) for row in table.rows
            )
        return "\n".join(chunks)[:max_chars], "local_docx"

    if lower_name.endswith((".xlsx", ".xlsm")):
        from openpyxl import load_workbook

        workbook = load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
        chunks: list[str] = []
        total = 0
        try:
            for sheet in workbook.worksheets:
                chunks.append(f"Foglio: {sheet.title}")
                for row in sheet.iter_rows(values_only=True):
                    line = " | ".join(str(value) for value in row if value is not None)
                    if line:
                        chunks.append(line)
                        total += len(line)
                    if total >= max_chars:
                        break
                if total >= max_chars:
                    break
        finally:
            workbook.close()
        return "\n".join(chunks)[:max_chars], "local_xlsx"

    if lower_name.endswith(".pptx"):
        from pptx import Presentation

        presentation = Presentation(io.BytesIO(raw))
        chunks = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                text = getattr(shape, "text", "")
                if text:
                    chunks.append(text)
        return "\n".join(chunks)[:max_chars], "local_pptx"

    return "", "unsupported"


def _adobe_pdf_to_markdown(raw: bytes) -> str:
    # Imported lazily so FinancePlus remains operational when Adobe credentials
    # are not configured or the optional SDK is not needed for a given file.
    from adobe.pdfservices.operation.auth.service_principal_credentials import (
        ServicePrincipalCredentials,
    )
    from adobe.pdfservices.operation.io.cloud_asset import CloudAsset
    from adobe.pdfservices.operation.io.stream_asset import StreamAsset
    from adobe.pdfservices.operation.pdf_services import PDFServices
    from adobe.pdfservices.operation.pdf_services_media_type import PDFServicesMediaType
    from adobe.pdfservices.operation.pdfjobs.jobs.pdf_to_markdown_job import (
        PDFToMarkdownJob,
    )
    from adobe.pdfservices.operation.pdfjobs.result.pdf_to_markdown_result import (
        PDFToMarkdownResult,
    )

    client_id = os.getenv("PDF_SERVICES_CLIENT_ID", "").strip()
    client_secret = os.getenv("PDF_SERVICES_CLIENT_SECRET", "").strip()
    if not client_id or not client_secret:
        raise RuntimeError("Credenziali Adobe PDF Services non configurate")

    credentials = ServicePrincipalCredentials(
        client_id=client_id, client_secret=client_secret
    )
    pdf_services = PDFServices(credentials=credentials)
    input_asset = pdf_services.upload(
        input_stream=raw, mime_type=PDFServicesMediaType.PDF
    )
    job = PDFToMarkdownJob(input_asset=input_asset)
    location = pdf_services.submit(job)
    response = pdf_services.get_job_result(location, PDFToMarkdownResult)
    result_asset: CloudAsset = response.get_result().get_asset()
    stream_asset: StreamAsset = pdf_services.get_content(result_asset)
    content = stream_asset.get_input_stream()
    if isinstance(content, bytes):
        return content.decode("utf-8", errors="replace").strip()
    return bytes(content).decode("utf-8", errors="replace").strip()


def _cloud_allowed(sensitivity: str, ai_policy: str) -> bool:
    if ai_policy == "Bloccata" or sensitivity == "Altamente riservato":
        return False
    if sensitivity == "Riservato":
        return _truthy("FINANCEPLUS_ADOBE_ALLOW_CONFIDENTIAL", default=False)
    return True


def extract_document_content(
    raw: bytes,
    filename: str,
    mime_type: str = "",
    sensitivity: str = "Interno",
    ai_policy: str = "Consentita",
    allow_cloud: bool | None = None,
) -> ExtractionResult:
    """Extract document text without weakening FinancePlus privacy controls.

    PDFs are always eligible for local extraction. Adobe PDF-to-Markdown is an
    optional quality layer. Callers handling unclassified inbound documents
    should first call this function with ``allow_cloud=False``, classify the
    local result, then call it again with ``allow_cloud=True`` only if the
    resulting privacy policy permits cloud processing.
    """
    lower_name = (filename or "").lower()
    is_pdf = lower_name.endswith(".pdf") or mime_type == "application/pdf"
    is_image = lower_name.endswith(
        (".jpg", ".jpeg", ".png", ".tif", ".tiff", ".webp", ".heic", ".heif")
    ) or mime_type.startswith("image/")

    if not raw:
        return ExtractionResult(warnings=["Documento vuoto"])

    if not is_pdf:
        if lower_name.endswith(
            (".txt", ".csv", ".md", ".json", ".xml")
        ) or mime_type.startswith("text/"):
            return ExtractionResult(
                text=raw.decode("utf-8", errors="replace")[:250_000],
                method="local_text",
            )
        if is_image:
            try:
                text = _local_image_ocr(raw)
                return ExtractionResult(
                    text=text,
                    method="local_tesseract_image" if text else "none",
                    warnings=[]
                    if text
                    else ["OCR locale non ha rilevato testo nell'immagine"],
                )
            except Exception as exc:
                return ExtractionResult(
                    method="ocr_failed", warnings=[f"OCR immagine non riuscito: {exc}"]
                )
        try:
            text, method = _local_office_text(raw, lower_name)
            return ExtractionResult(text=text, method=method)
        except Exception as exc:
            return ExtractionResult(
                method="office_failed",
                warnings=[f"Estrazione documento Office non riuscita: {exc}"],
            )

    warnings: list[str] = []
    local_text = ""
    try:
        local_text = _local_pdf_text(raw)
    except Exception as exc:
        warnings.append(f"Estrazione PDF locale non riuscita: {exc}")

    local_method = "local_pypdf" if local_text else "none"
    if len(local_text.strip()) < 80:
        try:
            ocr_text = _local_scanned_pdf_ocr(raw)
            if len(ocr_text) > len(local_text):
                local_text = ocr_text
                local_method = "local_tesseract_pdf"
            elif not ocr_text and not local_text:
                warnings.append("OCR locale non ha rilevato testo nel PDF scansionato")
        except Exception as exc:
            warnings.append(f"OCR PDF locale non riuscito: {exc}")

    mode = os.getenv("FINANCEPLUS_PDF_EXTRACTOR", "auto").strip().lower() or "auto"
    adobe_configured = bool(
        os.getenv("PDF_SERVICES_CLIENT_ID", "").strip()
        and os.getenv("PDF_SERVICES_CLIENT_SECRET", "").strip()
    )
    cloud_requested = True if allow_cloud is None else bool(allow_cloud)

    should_try_adobe = (
        cloud_requested
        and mode in {"auto", "adobe"}
        and adobe_configured
        and _cloud_allowed(sensitivity, ai_policy)
    )
    if should_try_adobe:
        try:
            markdown = _adobe_pdf_to_markdown(raw)
            if markdown:
                return ExtractionResult(
                    text=markdown[:400_000],
                    method="adobe_pdf_to_markdown",
                    cloud_used=True,
                    warnings=warnings,
                )
        except Exception as exc:
            warnings.append(f"Adobe PDF to Markdown non riuscito: {exc}")
            if mode == "adobe" and not local_text:
                return ExtractionResult(method="adobe_failed", warnings=warnings)

    if (
        cloud_requested
        and mode == "adobe"
        and not _cloud_allowed(sensitivity, ai_policy)
    ):
        warnings.append(
            "Adobe bloccato dalla policy privacy FinancePlus; usata estrazione locale"
        )

    return ExtractionResult(
        text=local_text,
        method=local_method,
        cloud_used=False,
        warnings=warnings,
    )
